import io
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def parse_file(file):
    ext = file.filename.split('.')[-1]
    content = ""
    if ext == "pdf":
        reader = PdfReader(file)
        content = "\n".join(page.extract_text() for page in reader.pages)
    elif ext == "docx":
        doc = Document(file)
        content = "\n".join(p.text for p in doc.paragraphs)
    elif ext == "pptx":
        ppt = Presentation(file)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
    elif ext == "csv":
        df = pd.read_csv(file)
        content = df.to_string()
    elif ext in ["txt", "md"]:
        content = file.read().decode("utf-8")
    return content